import os
import subprocess
from pptx import Presentation
from PIL import Image


def convert_pptx_to_images(pptx_file, output_folder):
    prs = Presentation(pptx_file)
    for i, slide in enumerate(prs.slides):
        image_width = 800
        image_height = 800
        image = Image.new('RGB', (image_width, image_height), (255, 255, 255))
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                image_path = f"{C:\\Users\\nooka\\OneDrive\\Desktop\\Handgesture-Recognition\\HandGesture.pptx}/slide_{i + 1}.png"
                shape.image.save(image_path)
                break
        image.save(image_path)

pptx_file = "HandGesture.pptx"
output_folder = "PRESENTATION"
convert_pptx_to_images(pptx_file, output_folder)