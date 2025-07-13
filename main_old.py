# Import necessary libraries
import streamlit as st
import os
import subprocess
import aspose.slides as slides
import aspose.pydrawing as drawing
from pptx import Presentation
import sys
from pathlib import Path

# Add src directory to path for imports
sys.path.append(str(Path(__file__).parent / "src"))

def convert_ppt_to_png(pptx_path, output_folder):
    """Convert PowerPoint presentation to PNG images."""
    try:
        # Load presentation
        pres = slides.Presentation(pptx_path)

        # Create the output folder if it doesn't exist
        os.makedirs(output_folder, exist_ok=True)

        # Loop through slides
        for index in range(pres.slides.length):
            # Get reference to slide
            slide = pres.slides[index]

            # Save as image
            image_path = os.path.join(output_folder, f"{index + 1}.jpeg")
            slide.get_thumbnail().save(image_path, drawing.imaging.ImageFormat.jpeg)

        st.success(f"✅ Converted {pres.slides.length} slides to images")
        return True
    except Exception as e:
        st.error(f"❌ Error converting presentation: {e}")
        return False

def main():
    """Main Streamlit application."""
    
    # Set page config
    st.set_page_config(
        page_title="HandGesture Recognition",
        page_icon="🎤",
        layout="wide"
    )
    
    # Header
    st.title("🎤 HandGesture Recognition")
    st.subheader("Control Presentations with Hand Gestures")
    st.markdown("---")

    # Sidebar for navigation
    st.sidebar.title("Navigation")
    page = st.sidebar.selectbox(
        "Choose a page",
        ["Upload & Convert", "Gesture Control", "Settings"]
    )

    if page == "Upload & Convert":
        upload_and_convert_page()
    elif page == "Gesture Control":
        gesture_control_page()
    elif page == "Settings":
        settings_page()

def upload_and_convert_page():
    """Page for uploading and converting presentations."""
    st.header("📤 Upload & Convert Presentation")
    
    # File upload for presentation
    uploaded_file = st.file_uploader(
        "Upload Presentation File (PPTX)", 
        type=["pptx"],
        help="Select a PowerPoint presentation file"
    )

    col1, col2 = st.columns(2)

    with col1:
        # Convert slides button
        if st.button("🔄 Convert Slides", type="primary"):
            if uploaded_file is not None:
                # Save the uploaded PPTX file
                pptx_path = "data/presentations/uploaded_presentation.pptx"
                os.makedirs("data/presentations", exist_ok=True)
                
                with open(pptx_path, "wb") as pptx_file:
                    pptx_file.write(uploaded_file.read())

                # Convert PPTX to images
                output_folder = "data/slides/images"
                if convert_ppt_to_png(pptx_path, output_folder):
                    st.success("✅ Slides converted successfully!")
                    st.info("💡 You can now use gesture control to navigate through the slides.")
            else:
                st.warning("⚠️ Please upload a presentation file first.")

    with col2:
        # Start gesture control button
        if st.button("🎮 Start Gesture Control"):
            if os.path.exists("data/slides/images") and len(os.listdir("data/slides/images")) > 0:
                st.info("🚀 Starting gesture controller...")
                
                # Execute gesture controller
                try:
                    process = subprocess.Popen(
                        ["python", "src/gesture.py"], 
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE
                    )
                    stdout, stderr = process.communicate()
                    
                    if process.returncode == 0:
                        st.success("✅ Gesture controller started successfully!")
                    else:
                        st.error(f"❌ Error starting gesture controller: {stderr.decode('utf-8')}")
                        
                except Exception as e:
                    st.error(f"❌ Failed to start gesture controller: {e}")
            else:
                st.warning("⚠️ Please convert slides first before starting gesture control.")

    # Display current slides
    if os.path.exists("data/slides/images"):
        slides_count = len([f for f in os.listdir("data/slides/images") 
                          if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
        if slides_count > 0:
            st.subheader(f"📊 Current Slides ({slides_count} total)")
            
            # Show first few slides as preview
            slide_files = sorted([f for f in os.listdir("data/slides/images") 
                                if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
            
            cols = st.columns(min(3, len(slide_files)))
            for i, slide_file in enumerate(slide_files[:3]):
                with cols[i]:
                    slide_path = os.path.join("data/slides/images", slide_file)
                    st.image(slide_path, caption=f"Slide {i+1}", use_column_width=True)

def gesture_control_page():
    """Page for gesture control information."""
    st.header("🎮 Gesture Control Guide")
    
    st.markdown("""
    ### Hand Gestures
    
    Position your hand at face level (above the green threshold line) and use these gestures:
    
    | Gesture | Action | Finger Pattern |
    |---------|--------|----------------|
    | 👆 Index finger only | Previous slide / Draw | `[1, 0, 0, 0, 0]` |
    | 👍 Thumb up | Next slide | `[0, 0, 0, 0, 1]` |
    | 👆 Index finger | Draw/Annotate | `[0, 1, 0, 0, 0]` |
    | 🤟 Three fingers | Erase annotation | `[0, 1, 1, 1, 0]` |
    | ✌️ Two fingers | Pointer mode | `[0, 1, 1, 0, 0]` |
    
    ### Keyboard Shortcuts
    
    - `q`: Quit the application
    - `r`: Reset all annotations
    - `n`: Next slide
    - `p`: Previous slide
    
    ### Tips for Best Results
    
    1. **Lighting**: Ensure good, consistent lighting
    2. **Background**: Use plain backgrounds
    3. **Distance**: Keep hand 1-2 feet from camera
    4. **Stability**: Minimize camera movement
    5. **Position**: Keep hand above the green threshold line
    """)

def settings_page():
    """Page for configuration settings."""
    st.header("⚙️ Settings")
    
    st.subheader("Configuration")
    
    if os.path.exists("config/gesture_config.json"):
        with open("config/gesture_config.json", "r") as f:
            import json
            config = json.load(f)
        
        st.json(config)
        
        if st.button("📝 Edit Configuration"):
            st.info("💡 Edit the config/gesture_config.json file to customize settings.")
    else:
        st.warning("⚠️ Configuration file not found.")
    
    st.subheader("Project Structure")
    st.code("""
Handgesture-Recognition/
├── src/                    # Source code
├── config/                 # Configuration files
├── data/                   # Data and assets
│   ├── presentations/      # Original PPTX files
│   ├── slides/            # Converted images
│   └── samples/           # Sample images
├── docs/                   # Documentation
└── assets/                 # Static assets
    """)

if __name__ == "__main__":
    main() 